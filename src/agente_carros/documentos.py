"""Leitura de documentos em varios formatos, para a indexacao.

O acervo do projeto e de PDF, mas uma base de conhecimento corporativa
chega em formato misto: politica em Word, tabela de precos em Excel,
apresentacao em PowerPoint, FAQ em Markdown ou HTML, exportacao de sistema
em CSV ou JSON. Cada formato exige um tratamento diferente para virar
texto que faca sentido depois de fatiado.

Duas decisoes que valem explicacao:

- As bibliotecas de cada formato sao importadas dentro da funcao que as
  usa. Quem so tem PDF no acervo nao precisa instalar nada a mais, e a
  ausencia de uma delas vira aviso no formato que faltou, nao erro na
  inicializacao do projeto.
- Planilha vira uma frase por linha, com o cabecalho repetido em cada
  celula. "Modelo: Corolla | Preco: 145000" sobrevive ao fatiamento;
  uma tabela em colunas perde o cabecalho no primeiro corte e o trecho
  recuperado fica com numeros soltos, sem dizer do que sao.
"""

from __future__ import annotations

import csv
import json
from pathlib import Path

FORMATOS = {
    ".pdf": "PDF",
    ".docx": "Word",
    ".xlsx": "Excel",
    ".pptx": "PowerPoint",
    ".md": "Markdown",
    ".markdown": "Markdown",
    ".txt": "Texto",
    ".csv": "CSV",
    ".json": "JSON",
    ".html": "HTML",
    ".htm": "HTML",
}

LIMITE_LINHAS_TABELA = 5000


class DependenciaAusente(RuntimeError):
    """Falta a biblioteca que le este formato."""


def formatos_suportados() -> list[str]:
    return sorted(FORMATOS)


def eh_suportado(caminho: Path) -> bool:
    return caminho.suffix.lower() in FORMATOS


def _linha_de_tabela(cabecalho: list[str], valores: list[str]) -> str:
    """Uma linha de planilha vira texto que se sustenta sozinho."""
    partes = [
        f"{titulo.strip()}: {str(valor).strip()}"
        for titulo, valor in zip(cabecalho, valores, strict=False)
        if str(valor).strip()
    ]
    return " | ".join(partes)


def ler_docx(caminho: Path) -> str:
    try:
        import docx  # type: ignore[import-untyped]
    except ImportError as erro:  # pragma: no cover - depende do ambiente
        raise DependenciaAusente("instale python-docx para ler .docx") from erro

    documento = docx.Document(str(caminho))
    partes = [p.text.strip() for p in documento.paragraphs if p.text.strip()]
    for tabela in documento.tables:
        linhas = [[celula.text for celula in linha.cells] for linha in tabela.rows]
        if not linhas:
            continue
        cabecalho = linhas[0]
        partes.extend(
            texto for linha in linhas[1:] if (texto := _linha_de_tabela(cabecalho, linha))
        )
    return "\n".join(partes)


def ler_xlsx(caminho: Path) -> str:
    try:
        from openpyxl import load_workbook  # type: ignore[import-untyped]
    except ImportError as erro:  # pragma: no cover - depende do ambiente
        raise DependenciaAusente("instale openpyxl para ler .xlsx") from erro

    livro = load_workbook(str(caminho), read_only=True, data_only=True)
    partes: list[str] = []
    for planilha in livro.worksheets:
        partes.append(f"# Planilha: {planilha.title}")
        cabecalho: list[str] = []
        for numero, linha in enumerate(planilha.iter_rows(values_only=True)):
            valores = ["" if v is None else str(v) for v in linha]
            if not any(v.strip() for v in valores):
                continue
            if not cabecalho:
                cabecalho = valores
                continue
            if numero > LIMITE_LINHAS_TABELA:
                partes.append(f"[planilha truncada em {LIMITE_LINHAS_TABELA} linhas]")
                break
            if texto := _linha_de_tabela(cabecalho, valores):
                partes.append(texto)
    livro.close()
    return "\n".join(partes)


def ler_pptx(caminho: Path) -> str:
    try:
        from pptx import Presentation  # type: ignore[import-untyped]
    except ImportError as erro:  # pragma: no cover - depende do ambiente
        raise DependenciaAusente("instale python-pptx para ler .pptx") from erro

    apresentacao = Presentation(str(caminho))
    partes: list[str] = []
    for numero, slide in enumerate(apresentacao.slides, start=1):
        partes.append(f"# Slide {numero}")
        for forma in slide.shapes:
            if getattr(forma, "has_text_frame", False) and forma.text_frame.text.strip():
                partes.append(forma.text_frame.text.strip())
        # As notas do apresentador costumam trazer o contexto que o slide
        # omite, e por isso entram no indice junto com ele.
        if slide.has_notes_slide:
            notas = slide.notes_slide.notes_text_frame.text.strip()
            if notas:
                partes.append(f"Notas do apresentador: {notas}")
    return "\n".join(partes)


def ler_html(caminho: Path) -> str:
    texto = caminho.read_text(encoding="utf-8", errors="ignore")
    try:
        from bs4 import BeautifulSoup  # type: ignore[import-untyped]
    except ImportError as erro:  # pragma: no cover - depende do ambiente
        raise DependenciaAusente("instale beautifulsoup4 para ler .html") from erro

    # html.parser e o do proprio Python: evita a dependencia de lxml.
    sopa = BeautifulSoup(texto, "html.parser")
    for indesejado in sopa(["script", "style", "nav", "footer"]):
        indesejado.decompose()
    return "\n".join(linha.strip() for linha in sopa.get_text("\n").splitlines() if linha.strip())


def ler_csv(caminho: Path) -> str:
    with caminho.open(encoding="utf-8", errors="ignore", newline="") as arquivo:
        amostra = arquivo.read(4096)
        arquivo.seek(0)
        try:
            dialeto = csv.Sniffer().sniff(amostra, delimiters=",;\t|")
        except csv.Error:
            dialeto = csv.excel
        leitor = csv.reader(arquivo, dialeto)
        linhas = []
        cabecalho: list[str] = []
        for numero, valores in enumerate(leitor):
            if not any(v.strip() for v in valores):
                continue
            if not cabecalho:
                cabecalho = valores
                continue
            if numero > LIMITE_LINHAS_TABELA:
                linhas.append(f"[arquivo truncado em {LIMITE_LINHAS_TABELA} linhas]")
                break
            if texto := _linha_de_tabela(cabecalho, valores):
                linhas.append(texto)
    return "\n".join(linhas)


def _achatar_json(valor, prefixo: str = "") -> list[str]:
    """Transforma o JSON em `caminho: valor`, um por linha.

    Manter o caminho ate a folha preserva o significado do dado depois do
    fatiamento: "beneficios.vale_refeicao.valor: 40" continua legivel,
    enquanto um "40" solto nao diz nada.
    """
    if isinstance(valor, dict):
        linhas: list[str] = []
        for chave, sub in valor.items():
            linhas.extend(_achatar_json(sub, f"{prefixo}.{chave}" if prefixo else str(chave)))
        return linhas
    if isinstance(valor, list):
        linhas = []
        for indice, sub in enumerate(valor):
            linhas.extend(_achatar_json(sub, f"{prefixo}[{indice}]"))
        return linhas
    return [f"{prefixo}: {valor}"] if prefixo else [str(valor)]


def ler_json(caminho: Path) -> str:
    dados = json.loads(caminho.read_text(encoding="utf-8", errors="ignore"))
    return "\n".join(_achatar_json(dados))


def ler_texto(caminho: Path) -> str:
    """Markdown e texto puro. O fatiador ja respeita a estrutura de titulos."""
    return caminho.read_text(encoding="utf-8", errors="ignore").strip()


LEITORES = {
    ".docx": ler_docx,
    ".xlsx": ler_xlsx,
    ".pptx": ler_pptx,
    ".html": ler_html,
    ".htm": ler_html,
    ".csv": ler_csv,
    ".json": ler_json,
    ".md": ler_texto,
    ".markdown": ler_texto,
    ".txt": ler_texto,
}


def extrair_texto(caminho: Path) -> str:
    """Texto de um documento nao-PDF. PDF continua com o PyPDFLoader.

    Nao devolve paginas porque so o PDF as tem; a paginacao dos demais
    formatos e uma invencao do editor, nao do documento.
    """
    extensao = caminho.suffix.lower()
    leitor = LEITORES.get(extensao)
    if leitor is None:
        raise ValueError(f"Formato nao suportado: {extensao}")
    return leitor(caminho)
